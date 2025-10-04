from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()

slides = [
    {
        'title': 'EventEye — Certificate Management System',
        'bullets': [
            'PDF + QR certificate generation, verification & admin UI',
            'Demo goal: generate → verify → manage (filter/sort) → delete/restore'
        ],
        'notes': 'Elevator pitch: A complete platform for issuing and managing certificates for events and courses.'
    },
    {
        'title': 'Architecture',
        'bullets': [
            'Backend: Python 3.11 + Flask',
            'Database: SQLite (local file)',
            'Frontend: Bootstrap 5 + Jinja2 + JavaScript'
        ],
        'notes': 'Explain the simple, self-contained architecture: easy to deploy, easy to extend. Point to files: enhanced_web.py, certificate_generator.py, templates/'
    },
    {
        'title': 'Core Features',
        'bullets': [
            'Single certificate generation (PDF + QR)',
            'Bulk generation (CSV)',
            'Verification via QR or ID',
            'Admin UI: search, filter, sort',
            'Soft delete + restore',
            'Dark/Light theme'
        ],
        'notes': 'Describe each feature briefly; show where to find them in the UI.'
    },
    {
        'title': 'Demo Plan',
        'bullets': [
            'Start server',
            'Generate a certificate',
            'Open and verify (QR / ID)',
            'Filter & sort in admin UI',
            'Soft delete & restore'
        ],
        'notes': 'Use the Live demo checklist provided. Keep steps short and clear.'
    },
    {
        'title': 'Tools & Frameworks',
        'bullets': [
            'Python 3.11, Flask',
            'ReportLab, Pillow, qrcode',
            'SQLite, Bootstrap 5, Font Awesome',
            'python-pptx (for this presentation)'
        ],
        'notes': 'Mention where each tool is used (backend, PDF, images, frontend).' 
    },
    {
        'title': 'Key Functions / Files',
        'bullets': [
            'enhanced_web.py — Flask routes & API',
            'certificate_generator.py — PDF/QR + DB CRUD',
            'templates/ — UI (base.html, certificates_list.html, etc.)',
            'certificates.db — data store'
        ],
        'notes': 'Map functions to files with one-line descriptions.'
    },
    {
        'title': 'Data Model',
        'bullets': [
            'Table: certificates',
            'Important fields: certificate_id, recipient_name, course_name, issue_date, organization, grade, certificate_hash, is_active'
        ],
        'notes': 'Explain soft delete via is_active and why it is useful.'
    },
    {
        'title': 'Security & Verification',
        'bullets': [
            'SHA-256 hash for tamper detection',
            'Unique certificate IDs',
            'Parameterized SQL + Jinja2 escaping'
        ],
        'notes': 'Shortly explain verification flow: QR contains verification URL + data; server checks hash and is_active.'
    },
    {
        'title': 'UX Highlights',
        'bullets': [
            'Real-time search, filter (date/grade/org) and sort',
            'Dark/Light theme with persistence',
            'Toasts and modals for feedback'
        ],
        'notes': 'Mention how filters and sorting improve admin productivity.'
    },
    {
        'title': 'Live Demo Commands',
        'bullets': [
            'Start: python enhanced_web.py',
            'Open: http://localhost:5000',
            'Create test: use Generate form or run a small script'
        ],
        'notes': 'Show the exact PowerShell commands in the demo.'
    },
    {
        'title': 'Next Steps / Roadmap',
        'bullets': [
            'Email distribution, Authentication, API endpoints',
            'Database migration to PostgreSQL, Background workers (Celery), Dockerize'
        ],
        'notes': 'Keep this slide to talk about future work and scalability.'
    },
    {
        'title': 'Q&A',
        'bullets': [
            'Be ready to explain PDF generation, verification, and scaling.'
        ],
        'notes': 'Have a couple of answers ready for common questions.'
    }
]

for s in slides:
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = s['title']
    body = slide.shapes.placeholders[1].text_frame
    body.text = s['bullets'][0]
    for bullet in s['bullets'][1:]:
        p = body.add_paragraph()
        p.text = bullet
        p.level = 1
    # Notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = s['notes']

out_path = os.path.join(os.getcwd(), 'EventEye_Presentation.pptx')
prs.save(out_path)
print('Presentation created at:', out_path)